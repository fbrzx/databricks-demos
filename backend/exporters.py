"""Turn a columns + rows result into downloadable CSV / Excel bytes."""
from __future__ import annotations

import csv
import io
from typing import Any
from xml.sax.saxutils import escape


def to_csv(columns: list[str], rows: list[list[Any]]) -> bytes:
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(columns)
    writer.writerows(rows)
    return buf.getvalue().encode("utf-8")


def to_xlsx(
    columns: list[str], rows: list[list[Any]], sheet_name: str = "Report"
) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name[:31] or "Report"

    ws.append(columns)
    for cell in ws[1]:
        cell.font = Font(bold=True)

    for row in rows:
        ws.append(["" if v is None else v for v in row])

    # Reasonable column widths based on header length.
    for i, name in enumerate(columns, start=1):
        ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = max(
            12, min(40, len(str(name)) + 4)
        )

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def to_pdf_report(report: dict[str, Any]) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        Paragraph,
        SimpleDocTemplate,
        Spacer,
    )

    primary = HexColor("#1A2536")
    base = HexColor("#F8F4EA")
    tan = HexColor("#B18A56")
    burgundy = HexColor("#6F2531")
    green = HexColor("#2F6A4C")
    line = HexColor("#DED3C1")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=letter,
        rightMargin=0.55 * inch,
        leftMargin=0.55 * inch,
        topMargin=0.75 * inch,
        bottomMargin=0.55 * inch,
        title=_text(report.get("title"), "Genie Report"),
    )

    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="ReportTitle",
            parent=styles["Title"],
            fontSize=22,
            leading=26,
            textColor=primary,
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SectionTitle",
            parent=styles["Heading2"],
            fontSize=12,
            leading=15,
            textColor=primary,
            spaceBefore=14,
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Body",
            parent=styles["BodyText"],
            fontSize=9.5,
            leading=13,
            textColor=primary,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Muted",
            parent=styles["BodyText"],
            fontSize=8.5,
            leading=11,
            textColor=colors.HexColor("#637083"),
        )
    )
    styles.add(
        ParagraphStyle(
            name="ReportCode",
            parent=styles["Code"],
            fontName="Courier",
            fontSize=7.5,
            leading=9,
            textColor=base,
            backColor=primary,
            borderPadding=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="TableCell",
            parent=styles["BodyText"],
            fontSize=7.5,
            leading=9,
            textColor=primary,
        )
    )

    story = [
        Paragraph(
            _paragraph(report.get("title"), "Genie Report"), styles["ReportTitle"]
        ),
    ]

    description = report.get("description")
    if description:
        story.append(Paragraph(_paragraph(description), styles["Muted"]))

    narrative = report.get("narrative") or report.get("text")
    if narrative:
        story.extend(
            [
                Spacer(1, 8),
                Paragraph("Narrative", styles["SectionTitle"]),
                Paragraph(_paragraph(narrative), styles["Body"]),
            ]
        )

    chart = _chart_drawing(report.get("chart"), burgundy, tan, green, primary)
    if chart:
        story.extend([Paragraph("Chart", styles["SectionTitle"]), chart])

    table = _preview_table(report, styles["TableCell"], primary, base, line)
    if table:
        story.extend([Paragraph("Table Excerpt", styles["SectionTitle"]), table])

    sql = report.get("sql")
    if sql:
        story.extend(
            [
                Paragraph("Generated SQL", styles["SectionTitle"]),
                Paragraph(_paragraph(_text(sql)[:2400]), styles["ReportCode"]),
            ]
        )

    def draw_brand(canvas, document):
        width, height = document.pagesize
        canvas.saveState()
        canvas.setFillColor(primary)
        canvas.rect(0, height - 30, width, 30, stroke=0, fill=1)
        canvas.setFillColor(base)
        canvas.setFont("Helvetica-Bold", 9)
        canvas.drawString(document.leftMargin, height - 20, "Genie Reports")
        canvas.setFillColor(tan)
        canvas.rect(0, height - 32, width, 2, stroke=0, fill=1)
        canvas.setFillColor(colors.HexColor("#637083"))
        canvas.setFont("Helvetica", 7)
        canvas.drawRightString(
            width - document.rightMargin,
            18,
            f"Page {document.page}",
        )
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_brand, onLaterPages=draw_brand)
    return buf.getvalue()


def _preview_table(
    report: dict[str, Any],
    cell_style: Any,
    primary: Any,
    base: Any,
    line: Any,
) -> Any | None:
    from reportlab.platypus import Paragraph, Table, TableStyle

    table = report.get("table") or {}
    columns = table.get("columns") or report.get("columns") or []
    rows = table.get("rows") or (report.get("rows") or [])[:10]
    if not columns or not rows:
        return None

    visible_columns = columns[:6]
    visible_rows = rows[:18]
    data = [[Paragraph(_paragraph(column), cell_style) for column in visible_columns]]
    for row in visible_rows:
        data.append(
            [
                Paragraph(_paragraph(_truncate(_cell(row, index), 64)), cell_style)
                for index, _ in enumerate(visible_columns)
            ]
        )

    pdf_table = Table(
        data,
        colWidths=[520 / len(visible_columns)] * len(visible_columns),
        repeatRows=1,
        hAlign="LEFT",
    )
    pdf_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), base),
                ("TEXTCOLOR", (0, 0), (-1, 0), primary),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.35, line),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
            ]
        )
    )
    return pdf_table


def _chart_drawing(
    chart: dict[str, Any] | None,
    burgundy: Any,
    tan: Any,
    green: Any,
    primary: Any,
) -> Any | None:
    from reportlab.graphics.shapes import Drawing, Rect, String

    if not chart or not chart.get("data"):
        return None

    value_columns = chart.get("value_columns") or chart.get("valueColumns") or []
    if not value_columns:
        return None

    value_column = value_columns[0]
    data = chart["data"][:8]
    values = [_float(row.get(value_column)) for row in data]
    max_value = max(values) if values else 0
    if max_value <= 0:
        return None

    drawing = Drawing(480, 150)
    bar_left = 125
    bar_max = 280
    bar_height = 11
    row_gap = 17
    colors_for_rows = [burgundy, tan, green]

    for index, row in enumerate(data):
        y = 128 - (index * row_gap)
        value = _float(row.get(value_column))
        width = max(1, (value / max_value) * bar_max)
        label = _truncate(row.get("label"), 24)
        drawing.add(String(0, y + 1, label, fontSize=7.5, fillColor=primary))
        drawing.add(
            Rect(
                bar_left,
                y,
                width,
                bar_height,
                fillColor=colors_for_rows[index % len(colors_for_rows)],
                strokeColor=None,
            )
        )
        drawing.add(
            String(
                bar_left + width + 5,
                y + 1,
                _truncate(value, 14),
                fontSize=7.5,
                fillColor=primary,
            )
        )

    return drawing


def _paragraph(value: Any, default: str = "") -> str:
    return escape(_text(value, default)).replace("\n", "<br/>")


def _text(value: Any, default: str = "") -> str:
    if value is None:
        return default
    return str(value)


def _truncate(value: Any, limit: int) -> str:
    text = _text(value)
    return text if len(text) <= limit else f"{text[: limit - 3]}..."


def _float(value: Any) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return 0


def _cell(row: list[Any], index: int) -> Any:
    return row[index] if index < len(row) else None


# ---- PPTX report export --------------------------------------------------- #

_PPTX_PRIMARY = "1A2536"
_PPTX_TAN = "B18A56"
_PPTX_BURGUNDY = "6F2531"
_PPTX_IVORY = "F8F4EA"
_PPTX_MUTED = "637083"
_PPTX_LINE = "DED3C1"


def _rgb(hex6: str):
    from pptx.util import Pt  # noqa: F401 - side-effect-free import guard
    from pptx.dml.color import RGBColor

    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def to_pptx_report(report: dict[str, Any]) -> bytes:
    """One report → one branded PPTX slide."""
    return to_pptx_bundle([report])


def to_pptx_bundle(reports: list[dict[str, Any]]) -> bytes:
    """Multiple reports → one branded PPTX deck, one slide per report."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for report in reports:
        slide = prs.slides.add_slide(blank_layout)
        W = prs.slide_width
        H = prs.slide_height

        # ── Header bar ───────────────────────────────────────────────────────
        header = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            0,
            0,
            W,
            Inches(0.55),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = _rgb(_PPTX_PRIMARY)
        header.line.fill.background()

        brand_tf = header.text_frame
        brand_tf.word_wrap = False
        brand_p = brand_tf.paragraphs[0]
        brand_r = brand_p.add_run()
        brand_r.text = "Genie Reports"
        brand_r.font.size = Pt(11)
        brand_r.font.bold = True
        brand_r.font.color.rgb = _rgb(_PPTX_IVORY)

        # ── Tan accent line below header ─────────────────────────────────────
        accent = slide.shapes.add_shape(1, 0, Inches(0.55), W, Emu(18000))
        accent.fill.solid()
        accent.fill.fore_color.rgb = _rgb(_PPTX_TAN)
        accent.line.fill.background()

        # ── Slide background ─────────────────────────────────────────────────
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(_PPTX_IVORY)

        # ── Title ────────────────────────────────────────────────────────────
        title_text = _text(report.get("title"), "Genie Report")
        title_box = slide.shapes.add_textbox(
            Inches(0.55), Inches(0.7), Inches(8), Inches(0.75)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title_text
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = _rgb(_PPTX_PRIMARY)

        # ── Description ──────────────────────────────────────────────────────
        desc = _text(report.get("description"))
        if desc:
            desc_box = slide.shapes.add_textbox(
                Inches(0.55), Inches(1.4), Inches(8), Inches(0.4)
            )
            tf2 = desc_box.text_frame
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = desc
            r2.font.size = Pt(10)
            r2.font.color.rgb = _rgb(_PPTX_MUTED)

        # ── Narrative ────────────────────────────────────────────────────────
        narrative = _text(report.get("narrative") or report.get("text"))
        y_narrative = Inches(1.75)
        if narrative:
            narr_box = slide.shapes.add_textbox(
                Inches(0.55), y_narrative, Inches(5.8), Inches(2.2)
            )
            tf3 = narr_box.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            r3 = p3.add_run()
            r3.text = narrative[:600] + ("…" if len(narrative) > 600 else "")
            r3.font.size = Pt(9.5)
            r3.font.color.rgb = _rgb(_PPTX_PRIMARY)

        # ── Table excerpt on the right ───────────────────────────────────────
        tbl_data = report.get("table") or {}
        columns = tbl_data.get("columns") or report.get("columns") or []
        rows = (tbl_data.get("rows") or report.get("rows") or [])[:10]
        visible_cols = columns[:5]

        if visible_cols and rows:
            table_x = Inches(7.0)
            table_y = Inches(0.75)
            table_w = Inches(5.9)
            tbl = slide.shapes.add_table(
                len(rows) + 1,
                len(visible_cols),
                table_x,
                table_y,
                table_w,
                Inches(min(3.2, 0.32 * (len(rows) + 1))),
            ).table

            for c_idx, col_name in enumerate(visible_cols):
                cell = tbl.cell(0, c_idx)
                cell.text = str(col_name)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(_PPTX_PRIMARY)
                tf_h = cell.text_frame
                tf_h.paragraphs[0].runs[0].font.color.rgb = _rgb(_PPTX_IVORY)
                tf_h.paragraphs[0].runs[0].font.size = Pt(7.5)
                tf_h.paragraphs[0].runs[0].font.bold = True

            for r_idx, row in enumerate(rows):
                for c_idx in range(len(visible_cols)):
                    cell = tbl.cell(r_idx + 1, c_idx)
                    val = _cell(row, c_idx)
                    cell.text = _truncate(val, 28)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(
                        "FFFDF8" if r_idx % 2 == 0 else _PPTX_IVORY
                    )
                    tf_c = cell.text_frame
                    tf_c.paragraphs[0].runs[0].font.size = Pt(7)
                    tf_c.paragraphs[0].runs[0].font.color.rgb = _rgb(_PPTX_PRIMARY)

        # ── SQL (optional, small) ────────────────────────────────────────────
        sql = _text(report.get("sql"))
        if sql and len(sql) < 800:
            sql_y = Inches(5.9)
            sql_box = slide.shapes.add_textbox(
                Inches(0.55), sql_y, Inches(6), Inches(1.0)
            )
            tf_sql = sql_box.text_frame
            tf_sql.word_wrap = True
            p_sql = tf_sql.paragraphs[0]
            r_sql = p_sql.add_run()
            r_sql.text = sql[:400] + ("…" if len(sql) > 400 else "")
            r_sql.font.size = Pt(6.5)
            r_sql.font.color.rgb = _rgb(_PPTX_MUTED)

        # ── Footer ───────────────────────────────────────────────────────────
        footer = slide.shapes.add_shape(1, 0, H - Inches(0.28), W, Inches(0.28))
        footer.fill.solid()
        footer.fill.fore_color.rgb = _rgb(_PPTX_PRIMARY)
        footer.line.fill.background()

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def to_pdf_bundle(reports: list[dict[str, Any]]) -> bytes:
    """Multiple reports → one branded multi-page PDF."""
    from reportlab.lib import colors
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    primary = HexColor("#1A2536")
    base = HexColor("#F8F4EA")
    tan = HexColor("#B18A56")
    burgundy = HexColor("#6F2531")
    green = HexColor("#2F6A4C")
    line = HexColor("#DED3C1")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=letter,
        rightMargin=0.55 * inch,
        leftMargin=0.55 * inch,
        topMargin=0.75 * inch,
        bottomMargin=0.55 * inch,
        title="Genie Report Bundle",
    )

    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            "ReportTitle",
            parent=styles["Title"],
            fontSize=22,
            leading=26,
            textColor=primary,
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            "SectionTitle",
            parent=styles["Heading2"],
            fontSize=12,
            leading=15,
            textColor=primary,
            spaceBefore=14,
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            "Body",
            parent=styles["BodyText"],
            fontSize=9.5,
            leading=13,
            textColor=primary,
        )
    )
    styles.add(
        ParagraphStyle(
            "Muted",
            parent=styles["BodyText"],
            fontSize=8.5,
            leading=11,
            textColor=colors.HexColor("#637083"),
        )
    )
    styles.add(
        ParagraphStyle(
            "ReportCode",
            parent=styles["Code"],
            fontName="Courier",
            fontSize=7.5,
            leading=9,
            textColor=base,
            backColor=primary,
            borderPadding=8,
        )
    )
    styles.add(
        ParagraphStyle(
            "TableCell",
            parent=styles["BodyText"],
            fontSize=7.5,
            leading=9,
            textColor=primary,
        )
    )

    def draw_brand(canvas, document):
        w, h = document.pagesize
        canvas.saveState()
        canvas.setFillColor(primary)
        canvas.rect(0, h - 30, w, 30, stroke=0, fill=1)
        canvas.setFillColor(base)
        canvas.setFont("Helvetica-Bold", 9)
        canvas.drawString(document.leftMargin, h - 20, "Genie Reports")
        canvas.setFillColor(tan)
        canvas.rect(0, h - 32, w, 2, stroke=0, fill=1)
        canvas.setFillColor(colors.HexColor("#637083"))
        canvas.setFont("Helvetica", 7)
        canvas.drawRightString(w - document.rightMargin, 18, f"Page {document.page}")
        canvas.restoreState()

    story = []
    for idx, report in enumerate(reports):
        if idx > 0:
            story.append(PageBreak())
        story.append(
            Paragraph(
                _paragraph(report.get("title"), "Genie Report"), styles["ReportTitle"]
            )
        )
        if report.get("description"):
            story.append(Paragraph(_paragraph(report["description"]), styles["Muted"]))
        narrative = report.get("narrative") or report.get("text")
        if narrative:
            story.extend(
                [
                    Spacer(1, 8),
                    Paragraph("Narrative", styles["SectionTitle"]),
                    Paragraph(_paragraph(narrative), styles["Body"]),
                ]
            )
        chart = _chart_drawing(report.get("chart"), burgundy, tan, green, primary)
        if chart:
            story.extend([Paragraph("Chart", styles["SectionTitle"]), chart])
        tbl = _preview_table(report, styles["TableCell"], primary, base, line)
        if tbl:
            story.extend([Paragraph("Table Excerpt", styles["SectionTitle"]), tbl])
        if report.get("sql"):
            story.extend(
                [
                    Paragraph("Generated SQL", styles["SectionTitle"]),
                    Paragraph(
                        _paragraph(_text(report["sql"])[:2400]), styles["ReportCode"]
                    ),
                ]
            )

    doc.build(story, onFirstPage=draw_brand, onLaterPages=draw_brand)
    return buf.getvalue()
